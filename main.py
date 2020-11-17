import tkinter as tk
import os
from tkinter import Tk
from tkinter.filedialog import askopenfilename, asksaveasfilename

from docx.shared import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Pt

HEIGHT = 250
WIDTH = 450

root = tk.Tk()

root.geometry("+500+300")

canvas = tk.Canvas(root, height=HEIGHT, width=WIDTH)
canvas.pack()

file = tk.PhotoImage(file="final.png")
image = canvas.create_image(200, 150, image=file)

global document
global prs
global fileName
global filename

separateDocument = False

e = False


def askopen():
    return askopenfilename(defaultextension=".pptx")


def prsThing():
    return Presentation(filename)


def asksave():
    return asksaveasfilename(defaultextension=".docx")


def documentThing():
    return Document()


def pptxGet():
    global filename
    filename = askopen()
    global prs
    prs = prsThing()


def docxGet():
    global fileName
    fileName = asksave()
    global document
    document = documentThing()


def imageSetting():
    global separateDocument
    if checked.get():
        separateDocument = True
    else:
        separateDocument = False


def program():
    global document
    global prs
    global fileName
    global filename
    global separateDocument

    tablesAndCharts = document.add_paragraph("There is")

    if separateDocument:
        imageDocument = Document()
    slideNum = 1
    for slide in prs.slides:
        p = document.add_paragraph()
        r = p.add_run("Slide " + str(slideNum) + ":")
        r.font.size = Pt(13)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                if shape.has_chart:
                    tablesAndCharts.add_run(f" a chart at slide {slideNum},")
                elif shape.has_table:
                    tablesAndCharts.add_run(f" a table at slide {slideNum},")
                continue
            for paragraph in shape.text_frame.paragraphs:
                count = 0
                for run in paragraph.runs:
                    if count == 0:
                        p = document.add_paragraph()
                        p.paragraph_format.space_after = Pt(5)
                        p.paragraph_format.line_spacing = Pt(0)
                    r = p.add_run(run.text)
                    r.font.size = Pt(9)
                    count += 1
        if separateDocument:
            savePictures(imageDocument, slide, slideNum, separateDocument)
        else:
            savePictures(document, slide, slideNum, separateDocument)
        slideNum += 1
    if tablesAndCharts.text == "There is":
        tablesAndCharts.add_run(" no tables or charts in the presentation")
    else:
        tablesAndCharts.add_run(" please refer to the presentation")
    document.save(fileName)
    if separateDocument:
        imageFile = str(fileName).replace(".docx", "Images.docx")
        imageDocument.save(imageFile)

def savePictures(file, slide, num, boo):
    picDescription = file.add_paragraph()
    if boo:
        picDescrip = picDescription.add_run(f"Slide {num}:")
    else:
        picDescrip = picDescription.add_run("Pictures on this slide:")
    picDescrip.font.size = Pt(11)
    pictures = file.add_paragraph()
    pics = pictures.add_run()

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            pic = image.blob
            with open('imageFile.jpg', 'wb') as f:
                f.write(pic)
            pics.add_picture("imageFile.jpg", height=Inches(1.0))
            os.remove("imageFile.jpg")


frame = tk.Frame(root, bg='#53ecff', height=HEIGHT, width=WIDTH)
frame.place(relx=.25, rely=.02, relwidth=.7, relheight=.62)

checked = tk.BooleanVar()
checked.set(False)
checkBox3 = tk.Checkbutton(frame)
checkBox3.config(text="Images on a separate Word Document",
                 command=imageSetting, anchor=tk.W,
                 bg='#53ecff', var=checked)
checkBox3.place(relx=.03, rely=.04)

inputFileFirst = tk.Label(frame, text="1.", bg='#53ecff')
inputFileFirst.place(relx=0.03, rely=.28)

inputFile = tk.Button(frame, text="Click", command=pptxGet)
inputFile.place(relx=0.092, rely=.27, relwidth=.12)

inputFileLabel = tk.Label(frame, text="to select a .pptx file to convert", bg='#53ecff')
inputFileLabel.place(relx=0.21, rely=.285)

outputFileSecond = tk.Label(frame, text="2.", bg='#53ecff')
outputFileSecond.place(relx=0.03, rely=.54)

outputFile = tk.Button(frame, text="Click", command=docxGet)
outputFile.place(relx=0.092, rely=.53, relwidth=.12)

outputFileLabel = tk.Label(frame, text="to create or select a .docx file to save to", bg='#53ecff')
outputFileLabel.place(relx=0.21, rely=.554)

programRunnerThird = tk.Label(frame, text="3.", bg='#53ecff')
programRunnerThird.place(relx=0.03, rely=.79)

programRunner = tk.Button(frame, text="Click", command=program)
programRunner.place(relx=0.092, rely=.78, relwidth=.12)

programRunnerLabel = tk.Label(frame, text="to run the program", bg='#53ecff')
programRunnerLabel.place(relx=0.21, rely=.795)

quitButton = tk.Button(frame, text="exit", command=root.quit)
quitButton.place(relx=.895, rely=.79, relwidth=.1)

root.mainloop()
